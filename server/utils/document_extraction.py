"""
Document Extraction Utility
============================

Extracts text content from various document formats in memory (no disk I/O).
Supports: TXT, MD, CSV, DOCX, XLSX, PDF, PPTX.
"""

import base64
import csv
import io
import logging

logger = logging.getLogger(__name__)

# Maximum characters of extracted text to send to Claude
MAX_EXTRACTED_CHARS = 200_000

# Maximum rows per sheet for Excel files
MAX_EXCEL_ROWS_PER_SHEET = 10_000
MAX_EXCEL_SHEETS = 50

# MIME type classification
DOCUMENT_MIME_TYPES: dict[str, str] = {
    "text/plain": ".txt",
    "text/markdown": ".md",
    "text/csv": ".csv",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}

IMAGE_MIME_TYPES = {"image/jpeg", "image/png"}

ALL_ALLOWED_MIME_TYPES = IMAGE_MIME_TYPES | set(DOCUMENT_MIME_TYPES.keys())


def is_image(mime_type: str) -> bool:
    """Check if the MIME type is a supported image format."""
    return mime_type in IMAGE_MIME_TYPES


def is_document(mime_type: str) -> bool:
    """Check if the MIME type is a supported document format."""
    return mime_type in DOCUMENT_MIME_TYPES


class DocumentExtractionError(Exception):
    """Raised when text extraction from a document fails."""

    def __init__(self, filename: str, reason: str):
        self.filename = filename
        self.reason = reason
        super().__init__(f"Failed to read {filename}: {reason}")


def _truncate(text: str) -> str:
    """Truncate text if it exceeds the maximum character limit."""
    if len(text) > MAX_EXTRACTED_CHARS:
        omitted = len(text) - MAX_EXTRACTED_CHARS
        return text[:MAX_EXTRACTED_CHARS] + f"\n\n[... truncated, {omitted:,} characters omitted]"
    return text


def _extract_plain_text(data: bytes) -> str:
    """Extract text from plain text or markdown files."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")


def _extract_csv(data: bytes) -> str:
    """Extract text from CSV files, formatted as a readable table."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1")

    reader = csv.reader(io.StringIO(text))
    lines = []
    for i, row in enumerate(reader):
        lines.append(f"Row {i + 1}: {', '.join(row)}")
    return "\n".join(lines)


def _extract_docx(data: bytes) -> str:
    """Extract text from Word documents."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from Excel spreadsheets."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sections = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        if sheet_idx >= MAX_EXCEL_SHEETS:
            sections.append(f"\n[... {len(wb.sheetnames) - MAX_EXCEL_SHEETS} more sheets omitted]")
            break

        ws = wb[sheet_name]
        rows_text = [f"=== Sheet: {sheet_name} ==="]
        row_count = 0

        for row in ws.iter_rows(values_only=True):
            if row_count >= MAX_EXCEL_ROWS_PER_SHEET:
                rows_text.append(f"[... more rows omitted, limit {MAX_EXCEL_ROWS_PER_SHEET:,} rows/sheet]")
                break
            cells = [str(cell) if cell is not None else "" for cell in row]
            rows_text.append("\t".join(cells))
            row_count += 1

        sections.append("\n".join(rows_text))

    wb.close()
    return "\n\n".join(sections)


def _extract_pdf(data: bytes, filename: str) -> str:
    """Extract text from PDF files."""
    from PyPDF2 import PdfReader
    from PyPDF2.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(data))
    except PdfReadError as e:
        if "encrypt" in str(e).lower() or "password" in str(e).lower():
            raise DocumentExtractionError(filename, "PDF is password-protected")
        raise

    if reader.is_encrypted:
        raise DocumentExtractionError(filename, "PDF is password-protected")

    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text}")

    return "\n\n".join(pages)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PowerPoint presentations."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides_text = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))

    return "\n\n".join(slides_text)


def extract_text_from_document(base64_data: str, mime_type: str, filename: str) -> str:
    """
    Extract text content from a document file.

    Args:
        base64_data: Base64-encoded file content
        mime_type: MIME type of the document
        filename: Original filename (for error messages)

    Returns:
        Extracted text content, truncated if necessary

    Raises:
        DocumentExtractionError: If extraction fails
    """
    if mime_type not in DOCUMENT_MIME_TYPES:
        raise DocumentExtractionError(filename, f"unsupported document type: {mime_type}")

    try:
        data = base64.b64decode(base64_data)
    except Exception as e:
        raise DocumentExtractionError(filename, f"invalid base64 data: {e}")

    try:
        if mime_type in ("text/plain", "text/markdown"):
            text = _extract_plain_text(data)
        elif mime_type == "text/csv":
            text = _extract_csv(data)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = _extract_docx(data)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            text = _extract_xlsx(data)
        elif mime_type == "application/pdf":
            text = _extract_pdf(data, filename)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = _extract_pptx(data)
        else:
            raise DocumentExtractionError(filename, f"unsupported document type: {mime_type}")
    except DocumentExtractionError:
        raise
    except Exception as e:
        logger.warning(f"Document extraction failed for {filename}: {e}")
        raise DocumentExtractionError(
            filename, "file appears to be corrupt or in an unexpected format"
        )

    if not text or not text.strip():
        return f"[File {filename} is empty or contains no extractable text]"

    return _truncate(text)
